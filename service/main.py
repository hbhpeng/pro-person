import glob
import os
import re
import sys
import string
from pathlib import Path
import secrets
import shutil

import openai
from icrawler import ImageDownloader
from icrawler.builtin import BingImageCrawler
from pptx import Presentation


def main(topic, slide_length, api_key, username):

    unique_image_name = secrets.token_hex(8)

    def refresh_unique_image_name():
        global unique_image_name
        unique_image_name = secrets.token_hex(8)
        return

    class PrefixNameDownloader(ImageDownloader):

        def get_filename(self, task, default_ext):
            global unique_image_name
            filename = super(PrefixNameDownloader, self).get_filename(
                task, default_ext)
            return 'prefix_' + unique_image_name + filename

    def generate_ppt(topic, slide_length, api_key, username):
        folder = Path() / "data" / "caches" / username
        if os.path.isdir(folder):
            shutil.rmtree(folder)
        folder.mkdir(parents=True, exist_ok=True)
        root = Presentation("./ppt-themes/theme0.pptx")

        openai.api_key = api_key

        message = f"""Create an outline for a slideshow presentation on the topic of {topic} which is {slide_length}
        slides long. Make sure it is {slide_length} long. Make sure The content is written in the Chinese.

        You are allowed to use the following slide types:
        Title Slide - (Title, Subtitle)
        Content Slide - (Title, Content)
        Image Slide - (Title, Content, Image)
        Thanks Slide - (Title)

        Put this tag before the Title Slide: [L_TS]
        Put this tag before the Content Slide: [L_CS]
        Put this tag before the Image Slide: [L_IS]
        Put this tag before the Thanks Slide: [L_THS]

        Put "[SLIDEBREAK]" after each slide

        For example:
        [L_TS]
        [TITLE]Mental Health[/TITLE]

        [SLIDEBREAK]

        [L_CS]
        [TITLE]Mental Health Definition[/TITLE]
        [CONTENT]
        1. Definition: A person’s condition with regard to their psychological and emotional well-being
        2. Can impact one's physical health
        3. Stigmatized too often.
        [/CONTENT]

        [SLIDEBREAK]

        Put this tag before the Title: [TITLE]
        Put this tag after the Title: [/TITLE]
        Put this tag before the Subitle: [SUBTITLE]
        Put this tag after the Subtitle: [/SUBTITLE]
        Put this tag before the Content: [CONTENT]
        Put this tag after the Content: [/CONTENT]
        Put this tag before the Image: [IMAGE]
        Put this tag after the Image: [/IMAGE]

        Elaborate on the Content, provide as much information as possible.
        You put a [/CONTENT] at the end of the Content.
        Do not reply as if you are talking about the slideshow itself. (ex. "Include pictures here about...")
        Do not include any special characters (?, !, ., :, ) in the Title.
        Do not include any additional information in your response and stick to the format."""

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": message}
            ]
        )

        # """ Ref for slide types:
        # 0 -> title and subtitle
        # 1 -> title and content
        # 2 -> section header
        # 3 -> two content
        # 4 -> Comparison
        # 5 -> Title only
        # 6 -> Blank
        # 7 -> Content with caption
        # 8 -> Pic with caption
        # """

        def delete_all_slides():
            for i in range(len(root.slides) - 1, -1, -1):
                r_id = root.slides._sldIdLst[i].rId
                root.part.drop_rel(r_id)
                del root.slides._sldIdLst[i]

        def create_title_slide(title, subtitle):
            layout = root.slide_layouts[0]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = subtitle

        def create_section_header_slide(title):
            layout = root.slide_layouts[2]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title

        def create_title_and_content_slide(title, content):
            layout = root.slide_layouts[1]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

        def create_title_and_content_and_image_slide(title, content, image_query):
            global unique_image_name
            layout = root.slide_layouts[8]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[2].text = content
            refresh_unique_image_name()
            google_crawler = BingImageCrawler(downloader_cls=PrefixNameDownloader, storage={'root_dir': folder})
            google_crawler.crawl(keyword=image_query, max_num=1)
            # dir_path = os.path.dirname(os.path.realpath(__file__))
            file_name = glob.glob(f"{folder}/prefix_{unique_image_name}*")
            img_path = ''
            if len(file_name) > 0:
                img_path = file_name[0]
                with open(img_path, 'rb') as f:
                    header = f.read(4)
                    if header != b'RIFF':
                        slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                     slide.placeholders[1].width, slide.placeholders[1].height)
                os.remove(img_path)

        def find_text_in_between_tags(text, start_tag, end_tag):
            start_pos = text.find(start_tag)
            end_pos = text.find(end_tag)
            result = []
            while start_pos > -1 and end_pos > -1:
                text_between_tags = text[start_pos + len(start_tag):end_pos]
                result.append(text_between_tags)
                start_pos = text.find(start_tag, end_pos + len(end_tag))
                end_pos = text.find(end_tag, start_pos)
            res1 = "".join(result)
            res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
            if len(result) > 0:
                return res2
            else:
                return ""

        def search_for_slide_type(text):
            tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
            found_text = next((s for s in tags if s in text), None)
            return found_text

        def parse_response(reply):
            list_of_slides = reply.split("[SLIDEBREAK]")
            for slide in list_of_slides:
                slide_type = search_for_slide_type(slide)
                if slide_type == "[L_TS]":
                    create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                       find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
                elif slide_type == "[L_CS]":
                    create_title_and_content_slide(
                        "".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                        "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                          "[/CONTENT]")))
                elif slide_type == "[L_IS]":
                    create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                               "[/TITLE]")),
                                                             "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                               "[/CONTENT]")),
                                                             "".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                                                               "[/IMAGE]")))
                elif slide_type == "[L_THS]":
                    create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

        def find_title():
            return root.slides[0].shapes.title.text

        delete_all_slides()

        # print(response)

        parse_response(response['choices'][0]['message']['content'])

        # root.save(f"{find_title()}.pptx")
        root.save(str(folder) + f"/{find_title()}.pptx")

        # print("done")

        return rf"Done! {find_title()} is ready! You can find it at {os.getcwd()}/{find_title()}.pptx"

    generate_ppt(topic, slide_length, api_key, username)


if __name__ == '__main__':
    arg1 = str(sys.argv[1])
    arg2 = int(sys.argv[2])
    arg3 = str(sys.argv[3])
    arg4 = str(sys.argv[4])
    main(arg1, arg2, arg3, arg4)
